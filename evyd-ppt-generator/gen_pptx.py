"""
EVYD PPT Generator — Free Mode Renderer
=========================================
Renders a content.json into a native-PPTX presentation.
All slides are drawn from code — no external template dependency.

Usage:
  python gen_pptx.py content.json [--style evyd_blue] [--output path.pptx]

Styles live in styles/<name>.json — pick one or create your own.
"""

import json, sys, os, argparse, math, glob as globmod
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Slide dimensions (20″ × 11.25″, 16:9) ────────────────────────────────────
SW = 20.0
SH = 11.25

I = Inches
P = Pt

CHROME_TYPES = {'cover', 'agenda', 'section_divider', 'ending'}


# ─────────────────────────────────────────────────────────────────────────────
# Style loader
# ─────────────────────────────────────────────────────────────────────────────

STYLES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'styles')

def _rgb(v):
    """Accept RGBColor | [r,g,b] list | '#RRGGBB' / 'RRGGBB' string."""
    if isinstance(v, RGBColor): return v
    if isinstance(v, (list, tuple)): return RGBColor(*v)
    if isinstance(v, str):
        h = v.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    raise ValueError(f"Cannot parse color: {v!r}")

def _is_light(color):
    """Return True if an RGBColor has relative luminance > 0.4 (perceptually light)."""
    r, g, b = color[0] / 255, color[1] / 255, color[2] / 255
    lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
    return lum > 0.4

def list_styles():
    """Return sorted list of available style names."""
    return sorted(
        os.path.splitext(os.path.basename(f))[0]
        for f in globmod.glob(os.path.join(STYLES_DIR, '*.json'))
    )

def load_style(name):
    """Load a style JSON file and return a flat dict with RGBColor values."""
    path = os.path.join(STYLES_DIR, f'{name}.json')
    if not os.path.exists(path):
        avail = ', '.join(list_styles())
        print(f'Unknown style "{name}". Available: {avail}')
        sys.exit(1)

    with open(path, encoding='utf-8') as f:
        raw = json.load(f)

    st = {}
    st['font'] = raw.get('font', 'Aptos')

    for key in ['accent', 'accent2', 'navy', 'white', 'card', 'card_side',
                'card_num', 'text_dim', 'text_num', 'card_white', 'text_gray',
                'text_dark', 'line_gray']:
        st[key] = _rgb(raw[key])

    st['bg_content']  = _rgb(raw.get('bg_content', raw['card']))
    st['bg_blue']     = st['bg_content']
    st['bg_white']    = _rgb(raw.get('white'))
    st['accent_dark'] = st['navy']
    st['header_bg']   = st['navy']
    st['header_sep']  = st['line_gray']
    st['text_white']  = st['white']
    st['text_primary']= st['text_dim']
    st['text_muted']  = st['text_gray']
    st['text_mid']    = st['text_gray']
    st['text_blue_lt']= st['text_num']
    st['text_example']= st['text_gray']
    st['slide_num_dk']= st['line_gray']
    st['slide_num_lt']= _rgb('AAAAAA')
    st['warning_bg']  = st['navy']

    # ── Light-background detection ────────────────────────────────────────────
    # If bg_content is perceptually light, the "blue" pages need dark text
    # instead of light text.  Chrome slides (cover, ending, section_divider)
    # always render on st['navy'] which stays dark — they use st['chrome_text'].
    st['blue_is_light'] = _is_light(st['bg_content'])
    st['chrome_text']   = st['white']           # always light — for dark chrome bg
    st['chrome_muted']  = st['text_gray'] if st['blue_is_light'] else st['text_dim']
    if st['blue_is_light']:
        st['text_white']   = st['text_dark']    # titles on "blue" bg → dark
        st['text_primary'] = st['text_gray']    # body on "blue" bg → medium-dark
        st['text_muted']   = st['text_gray']
        st['text_dim']     = st['text_gray']    # dim text → readable on light bg
        st['slide_num_dk'] = st['text_gray']
        st['slide_num_lt'] = st['line_gray']

    st['motifs'] = raw.get('motifs', {})

    # ── v2 schema additions (all optional, defensive defaults) ────────────────
    st['version']       = raw.get('version', 1)
    st['category']      = raw.get('category', 'classic')
    st['vibe_tags']     = raw.get('vibe_tags', [])
    st['chrome_style']  = raw.get('chrome_style', 'classic')
    st['title_font']    = raw.get('title_font', st['font'])
    st['body_font']     = raw.get('body_font',  st['font'])
    st['mono_font']     = raw.get('mono_font',  'Consolas')
    st['card_radius']   = raw.get('card_radius', True)
    st['gradient_spec'] = raw.get('gradient') or None
    st['motif_spec']    = raw.get('decorative_motif') or None

    # Chart series colors (fallback to accent, accent2, navy)
    cc = raw.get('chart_colors', [raw.get('accent', '2CD5C3'),
                                   raw.get('accent2', '0076B3'),
                                   raw.get('navy', '172E41')])
    st['chart_colors'] = [_rgb(c) for c in cc]

    return st


# ─────────────────────────────────────────────────────────────────────────────
# Primitive helpers
# ─────────────────────────────────────────────────────────────────────────────

def bx(slide, l, t, w, h, text='', sz=16, bold=False, color=None,
        align=None, italic=False, font='Aptos'):
    """Add a textbox. Returns the textbox shape."""
    tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    pg = tf.paragraphs[0]
    pg.alignment = align or PP_ALIGN.LEFT
    r = pg.add_run(); r.text = text; r.font.name = font
    r.font.size = P(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = _rgb(color)
    return tb

def ap(tf, text, sz=14, bold=False, color=None, align=None, spb=0,
        italic=False, font='Aptos'):
    """Append a paragraph to an existing text frame."""
    pg = tf.add_paragraph()
    pg.alignment = align or PP_ALIGN.LEFT
    if spb: pg.space_before = P(spb)
    r = pg.add_run(); r.text = text; r.font.name = font
    r.font.size = P(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = _rgb(color)

def rc(slide, l, t, w, h, fill=None, line=None, lw=0.75, rd=False):
    """Add a rectangle (rd=True for rounded). Returns the shape."""
    s = slide.shapes.add_shape(5 if rd else 1, I(l), I(t), I(w), I(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    else:    s.fill.background()
    if line: s.line.color.rgb = _rgb(line); s.line.width = P(lw)
    else:    s.line.fill.background()
    return s

def ov(slide, l, t, w, h, fill=None, transparency=0):
    """Add an oval shape with optional transparency."""
    s = slide.shapes.add_shape(9, I(l), I(t), I(w), I(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    if transparency:
        _set_transparency(s, transparency)
    s.line.fill.background()
    return s

def _set_transparency(shape, pct):
    """Set fill transparency (0-100) on a shape via XML."""
    from pptx.oxml.ns import qn
    from lxml import etree
    sp_pr = shape._element.find(qn('a:solidFill'))
    if sp_pr is None:
        sp_pr_elem = shape._element.find(qn('p:spPr'))
        if sp_pr_elem is None:
            return
        solid = sp_pr_elem.find(qn('a:solidFill'))
    else:
        solid = sp_pr
    if solid is None:
        return
    color_elem = solid[0] if len(solid) > 0 else None
    if color_elem is None:
        return
    alpha = color_elem.find(qn('a:alpha'))
    if alpha is None:
        alpha = etree.SubElement(color_elem, qn('a:alpha'))
    alpha.set('val', str(int((100 - pct) * 1000)))

def _fill_bg(slide, data, st):
    """Fill entire slide with the appropriate background color."""
    blue = data.get('background', 'blue') == 'blue'
    bg = st['bg_blue'] if blue else st['bg_white']
    rc(slide, 0, 0, SW, SH, fill=bg)


# ─────────────────────────────────────────────────────────────────────────────
# Shared header helper
# ─────────────────────────────────────────────────────────────────────────────

def hdr(slide, section, num, title, blue=True, tsz=34, st=None):
    """Standard content-slide header: section label + slide# + title + accent rule.
    Returns content_top_y (inches) — start your content below this.

    Hooks (v2 schema, optional):
      - st['motifs']['header_rule_color'] overrides the accent rule color.
      - st['title_font'] / st['body_font'] override the default font for title / label.
    """
    TF = (st.get('title_font') if st else None) or (st['font'] if st else 'Aptos')
    BF = (st.get('body_font')  if st else None) or TF
    lc = st['accent']   if st else _rgb('2CD5C3')
    nc = st['text_num'] if st else (_rgb('CCE8F5') if blue else _rgb('AAAAAA'))
    tc = st['text_white'] if st else (_rgb('FFFFFF') if blue else _rgb('172E41'))
    if not blue and st: tc = st['text_dark']
    if not blue and st: lc = st['accent2']
    # v2 hook: allow theme to override the header rule color independently.
    rule_c = lc
    if st:
        rule_override = st.get('motifs', {}).get('header_rule_color')
        if rule_override:
            rule_c = _rgb(rule_override)
    bx(slide, 1.0, 0.30, 10,   0.35, section, sz=11,  color=lc,  font=BF)
    bx(slide, 16.5, 0.30, 2.5, 0.35, num,     sz=11,  color=nc,  align=PP_ALIGN.RIGHT, font=BF)
    bx(slide, 1.0, 0.72, 17,   0.90, title,   sz=tsz, bold=True, color=tc, font=TF)
    rc(slide, 1.0, 1.60, 0.55, 0.055, fill=rule_c)
    return 1.85

def _slide_num(slide_data, auto_num, total):
    return slide_data.get('num', f'{auto_num:02d} / {total:02d}')

def _slide_num_free(slide, n, total, light=False, st=None):
    """Add slide number bottom-right."""
    color = st['slide_num_lt'] if light else st['slide_num_dk']
    bx(slide, 17.6, 10.6, 2, 0.4,
       f'{str(n).zfill(2)} / {str(total).zfill(2)}',
       sz=8 * 2, color=color, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# Chrome slide renderers — layered compositor
# ─────────────────────────────────────────────────────────────────────────────
#
# Architecture:
#   render_<chrome_slide>  →  _paint_chrome_bg   (solid or gradient)
#                          →  _paint_chrome_motif (optional, from motif_spec)
#                          →  _compose_<slot>_<style>  (typography + accents)
#                          →  _slide_num_free
#
# chrome_style values: classic | gradient | neon-grid | magazine | minimal | brutalist
# v1 themes (no chrome_style field) default to 'classic' — visually unchanged.

def _paint_gradient(slide, x, y, w, h, colors, angle, st, transparency=0):
    """Create a gradient-filled rectangle. Reuses freeform gradient logic."""
    from pptx.oxml.ns import qn
    from lxml import etree
    s = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    s.line.fill.background()
    gf = s.fill
    gf.gradient()
    gf.gradient_angle = angle
    stops = gf.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = _resolve_color(colors[0], st)
    stops[1].position = 1.0
    stops[1].color.rgb = _resolve_color(colors[-1], st)
    if len(colors) >= 3:
        gs_lst = s._element.find(qn('p:spPr')).find(qn('a:gradFill')).find(qn('a:gsLst'))
        for i, c in enumerate(colors[1:-1], 1):
            pos = int((i / (len(colors) - 1)) * 100000)
            gs = etree.SubElement(gs_lst, qn('a:gs'))
            gs.set('pos', str(pos))
            srgb = etree.SubElement(gs, qn('a:srgbClr'))
            rgb = _resolve_color(c, st)
            srgb.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')
    if transparency:
        _set_transparency(s, transparency)
    return s


def _paint_chrome_bg(slide, st, slot):
    """Paint full-slide background for a chrome slide.
    slot ∈ {'cover', 'section', 'ending'}.
    Uses gradient_spec[slot] if present, else solid navy."""
    gspec = st.get('gradient_spec')
    if gspec and isinstance(gspec, dict) and slot in gspec:
        g = gspec[slot]
        colors = g.get('colors', ['navy', 'accent'])
        angle  = g.get('angle', 0)
        _paint_gradient(slide, 0, 0, SW, SH, colors, angle, st)
    else:
        rc(slide, 0, 0, SW, SH, fill=st['navy'])


def _paint_chrome_motif(slide, st, slot):
    """Apply decorative motif to a chrome slide. No-op if motif_spec missing."""
    mspec = st.get('motif_spec')
    if not mspec:
        return
    kind = mspec.get('kind', 'none')
    if kind == 'none':
        return
    color   = _resolve_color(mspec.get('color', 'accent'), st)
    opacity = mspec.get('opacity', 85)
    density = mspec.get('density', 'medium')

    if   kind == 'oval':  _motif_oval(slide, slot, color, opacity)
    elif kind == 'grid':  _motif_grid(slide, slot, color, opacity, density)
    elif kind == 'blob':  _motif_blob(slide, slot, color, opacity)
    elif kind == 'lines': _motif_lines(slide, slot, color, opacity)
    elif kind == 'bars':  _motif_bars(slide, slot, color, opacity, st)
    elif kind == 'dots':  _motif_dots(slide, slot, color, opacity, density)


def _motif_oval(slide, slot, color, opacity):
    regions = {
        'cover':   [(11.6, -3.0, 12.0, 12.0)],
        'section': [(14.0, 6.0, 10.0, 10.0)],
        'ending':  [(5.0, 0.6, 10.0, 5.6)],
    }
    for (x, y, w, h) in regions.get(slot, []):
        s = slide.shapes.add_shape(9, I(x), I(y), I(w), I(h))
        s.fill.solid(); s.fill.fore_color.rgb = color
        _set_transparency(s, opacity)
        s.line.fill.background()


def _motif_grid(slide, slot, color, opacity, density):
    step = {'sparse': 2.0, 'medium': 1.2, 'dense': 0.8}.get(density, 1.2)
    line_w = 0.01
    y = step
    while y < SH:
        s = rc(slide, 0, y, SW, line_w, fill=color)
        _set_transparency(s, opacity)
        y += step
    x = step
    while x < SW:
        s = rc(slide, x, 0, line_w, SH, fill=color)
        _set_transparency(s, opacity)
        x += step


def _motif_blob(slide, slot, color, opacity):
    regions = {
        'cover':   [(15.0, 2.0, 6.0, 6.0), (1.0, 7.0, 4.0, 4.0)],
        'section': [(16.0, 1.0, 4.0, 4.0)],
        'ending':  [(2.0, 2.0, 5.0, 5.0), (14.0, 5.0, 5.0, 5.0)],
    }
    for (x, y, w, h) in regions.get(slot, []):
        s = slide.shapes.add_shape(9, I(x), I(y), I(w), I(h))
        s.fill.solid(); s.fill.fore_color.rgb = color
        _set_transparency(s, opacity)
        s.line.fill.background()


def _motif_lines(slide, slot, color, opacity):
    regions = {
        'cover':   [(16.5, 0.8, 0.02, 9.5)],
        'section': [(18.0, 0.8, 0.02, 9.5)],
        'ending':  [(1.2, 0.8, 0.02, 9.5), (18.8, 0.8, 0.02, 9.5)],
    }
    for (x, y, w, h) in regions.get(slot, []):
        s = rc(slide, x, y, w, h, fill=color)
        _set_transparency(s, max(0, opacity - 30))


def _motif_bars(slide, slot, color, opacity, st):
    accent2 = st['accent2']
    regions = {
        'cover':   [(0, 0, 0.14, SH / 2, color), (0, SH / 2, 0.14, SH / 2, accent2)],
        'section': [(0, 0, 0.14, SH, color)],
        'ending':  [(0, 0, 0.14, SH, color)],
    }
    for r in regions.get(slot, []):
        x, y, w, h, c = r
        rc(slide, x, y, w, h, fill=c)


def _motif_dots(slide, slot, color, opacity, density):
    import random
    random.seed(42 + hash(slot))
    n = {'sparse': 10, 'medium': 20, 'dense': 36}.get(density, 20)
    for _ in range(n):
        x = random.uniform(0.5, SW - 0.5)
        y = random.uniform(0.5, SH - 0.5)
        s = slide.shapes.add_shape(9, I(x), I(y), I(0.14), I(0.14))
        s.fill.solid(); s.fill.fore_color.rgb = color
        _set_transparency(s, opacity)
        s.line.fill.background()


# ── Cover composers ──────────────────────────────────────────────────────────

def _cover_classic(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    M = st['motifs']
    # Left bars (signature of classic chrome)
    if not st.get('motif_spec'):
        bar_colors = M.get('left_bar_colors', ['2CD5C3', '0076B3'])
        rc(slide, 0, 0, 0.14, SH / 2, fill=_rgb(bar_colors[0]))
        rc(slide, 0, SH / 2, 0.14, SH / 2, fill=_rgb(bar_colors[1]))
        ov(slide, 11.6, -3.0, 12, 12, fill=st['accent2'], transparency=94)

    tag_text = (data.get('tag', 'PRESENTATION')).upper()
    rc(slide, 1.1, 1.7, 5.6, 0.64, fill=st['accent_dark'],
       line=_rgb(M.get('header_tag_color', '2CD5C3')), lw=1)
    bx(slide, 1.1, 1.7, 5.6, 0.64, tag_text, sz=16, bold=True,
       color=_rgb(M.get('header_tag_color', '2CD5C3')),
       align=PP_ALIGN.CENTER, font=F)

    bx(slide, 1.1, 2.6, 14.4, 3.8,
       data.get('title', ''), sz=80, color=st['chrome_text'], font=F)

    rc(slide, 1.1, 6.7, 0.1, 1.36, fill=st['accent2'])
    bx(slide, 1.4, 6.7, 11.6, 1.36,
       data.get('subtitle', ''), sz=28, color=st['chrome_muted'], font=BF)

    bx(slide, 15.0, 10.2, 4.4, 0.6,
       data.get('logo', 'EVYD  ·  2025'), sz=18,
       color=st['header_sep'], align=PP_ALIGN.RIGHT, font=F)


def _cover_gradient(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    tag = (data.get('tag', '')).upper()
    if tag:
        bx(slide, 1.2, 1.4, 12, 0.5, tag, sz=14, bold=True,
           color=st['accent'], font=F)
        rc(slide, 1.2, 1.9, 1.0, 0.04, fill=st['accent'])
    bx(slide, 1.2, 3.0, 17.5, 4.2,
       data.get('title', ''), sz=88, bold=True,
       color=st['chrome_text'], font=F)
    bx(slide, 1.2, 7.2, 14.5, 1.5,
       data.get('subtitle', ''), sz=26,
       color=st['chrome_muted'], font=BF)
    bx(slide, 1.2, 10.2, 10, 0.6,
       data.get('logo', 'EVYD  ·  2025'), sz=14,
       color=st['chrome_muted'], font=F)


def _cover_neongrid(slide, data, st):
    F = st['title_font']; BF = st['body_font']; MF = st['mono_font']
    tag = (data.get('tag', '> READY')).upper()
    bx(slide, 1.1, 1.6, 12, 0.5, tag, sz=16, bold=True,
       color=st['accent'], font=MF)
    rc(slide, 1.1, 2.18, 0.8, 0.04, fill=st['accent'])
    bx(slide, 1.1, 2.9, 17.8, 4.5,
       data.get('title', ''), sz=80, bold=True,
       color=st['chrome_text'], font=F)
    bx(slide, 1.1, 7.6, 14, 1.4,
       data.get('subtitle', ''), sz=22,
       color=st['chrome_muted'], font=MF)
    bx(slide, 14.5, 10.2, 5, 0.6,
       data.get('logo', 'EVYD  ·  2025'), sz=12,
       color=st['accent2'], align=PP_ALIGN.RIGHT, font=MF)


def _cover_magazine(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    # Top hairline + issue label
    rc(slide, 1.5, 0.9, SW - 3.0, 0.02, fill=st['text_dim'])
    tag = (data.get('tag', 'ISSUE 01')).upper()
    bx(slide, 1.5, 1.05, 9, 0.4, tag, sz=12, bold=True,
       color=st['text_dim'], font=F)
    bx(slide, 10.5, 1.05, 8, 0.4,
       data.get('logo', 'EVYD · 2025'), sz=12,
       color=st['text_dim'], align=PP_ALIGN.RIGHT, font=F)
    # Serif title, centered
    bx(slide, 1.5, 3.4, SW - 3.0, 4.5,
       data.get('title', ''), sz=90, bold=True,
       color=st['chrome_text'], align=PP_ALIGN.CENTER, font=F)
    # Hairline divider
    rc(slide, 9.0, 7.9, 2.0, 0.03, fill=st['accent'])
    # Subtitle in italic
    bx(slide, 2.5, 8.3, SW - 5.0, 1.5,
       data.get('subtitle', ''), sz=22, italic=True,
       color=st['chrome_muted'], align=PP_ALIGN.CENTER, font=BF)
    rc(slide, 1.5, SH - 0.8, SW - 3.0, 0.02, fill=st['text_dim'])


def _cover_minimal(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    tag = (data.get('tag', '')).upper()
    if tag:
        bx(slide, 1.1, 1.2, 10, 0.4, tag, sz=11, bold=True,
           color=st['accent'], font=F)
    bx(slide, 1.1, 4.4, SW - 2.2, 3.5,
       data.get('title', ''), sz=76,
       color=st['chrome_text'], font=F)
    rc(slide, 1.1, 7.8, 0.6, 0.04, fill=st['accent'])
    bx(slide, 1.1, 8.1, 14, 1.2,
       data.get('subtitle', ''), sz=22,
       color=st['chrome_muted'], font=BF)
    bx(slide, 1.1, 10.2, 10, 0.6,
       data.get('logo', 'EVYD'), sz=11,
       color=st['chrome_muted'], font=F)


def _cover_brutalist(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    # Solid block for TAG
    tag = (data.get('tag', 'PRESENTATION')).upper()
    rc(slide, 1.1, 1.4, 5.5, 0.85, fill=st['accent'])
    bx(slide, 1.1, 1.5, 5.5, 0.7, tag, sz=18, bold=True,
       color=st['navy'], align=PP_ALIGN.CENTER, font=F)
    # Massive ALL-CAPS title
    bx(slide, 1.1, 2.8, SW - 2.2, 5.8,
       (data.get('title', '')).upper(), sz=100, bold=True,
       color=st['chrome_text'], font=F)
    # Heavy accent block
    rc(slide, 1.1, 8.8, 4.0, 0.25, fill=st['accent'])
    bx(slide, 1.1, 9.2, 14, 1.2,
       data.get('subtitle', ''), sz=24, bold=True,
       color=st['chrome_muted'], font=BF)
    bx(slide, 14.5, 10.2, 5, 0.6,
       data.get('logo', 'EVYD  ·  2025'), sz=13, bold=True,
       color=st['chrome_text'], align=PP_ALIGN.RIGHT, font=F)


_COVER_COMPOSERS = {
    'classic':   _cover_classic,
    'gradient':  _cover_gradient,
    'neon-grid': _cover_neongrid,
    'magazine':  _cover_magazine,
    'minimal':   _cover_minimal,
    'brutalist': _cover_brutalist,
}


def render_cover(slide, data, st, n, total):
    _paint_chrome_bg(slide, st, 'cover')
    _paint_chrome_motif(slide, st, 'cover')
    composer = _COVER_COMPOSERS.get(st.get('chrome_style', 'classic'),
                                    _cover_classic)
    composer(slide, data, st)
    _slide_num_free(slide, n, total, st=st)


def render_agenda(slide, data, st, _n, _total):
    F = st['font']
    M = st['motifs']
    rc(slide, 0, 0, SW, SH, fill=st['navy'])

    bx(slide, 1.0, 0.6, 10, 1.0, 'AGENDA', sz=28 * 2, bold=True,
       color=st['chrome_text'], font=F)
    rc(slide, 1.0, 1.6, 1.1, 0.09,
       fill=_rgb(M.get('divider_color', '0076B3')))

    items = data.get('items', [])
    for i, item in enumerate(items[:5]):
        y = 2.2 + i * 1.6
        rc(slide, 1.0, y, 18.0, 1.4, fill=st['card'])
        rc(slide, 1.0, y, 0.12, 1.4,
           fill=_rgb(M.get('number_color', '0076B3')))
        bx(slide, 1.3, y, 1.4, 1.4, str(item.get('num', i + 1)),
           sz=20 * 2, bold=True,
           color=_rgb(M.get('number_color', '0076B3')),
           align=PP_ALIGN.CENTER, font=F)
        rc(slide, 2.8, y + 0.35, 0.04, 0.7, fill=st['header_sep'])
        _card_text = st['text_dark'] if _is_light(st['card']) else st['chrome_text']
        _card_sub  = st['text_gray'] if _is_light(st['card']) else st['chrome_muted']
        bx(slide, 3.1, y + 0.1, 10, 0.7, item.get('title', ''),
           sz=15 * 2, bold=True, color=_card_text, font=F)
        bx(slide, 3.1, y + 0.75, 10, 0.5, item.get('time', ''),
           sz=11 * 2, color=_card_sub, font=F)


# ── Section divider composers ────────────────────────────────────────────────

def _section_classic(slide, data, st):
    F = st['title_font']
    M = st['motifs']
    bx(slide, 1.0, 2.5, 5, 3.5, str(data.get('num', '01')),
       sz=144, bold=True,
       color=_rgb(M.get('number_color', '0076B3')), font=F)
    rc(slide, 1.0, 5.8, 2.0, 0.08,
       fill=_rgb(M.get('divider_color', '0076B3')))
    bx(slide, 1.0, 6.2, 18, 2.0, data.get('title', 'Section'),
       sz=56, bold=True, color=st['chrome_text'], font=F)


def _section_gradient(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    num = str(data.get('num', '01'))
    bx(slide, 1.0, 1.6, 6, 2.5, num, sz=180, bold=True,
       color=st['accent'], font=F)
    rc(slide, 1.0, 6.0, 3.0, 0.06, fill=st['accent'])
    bx(slide, 1.0, 6.4, 18, 2.2, data.get('title', 'Section'),
       sz=60, bold=True, color=st['chrome_text'], font=F)
    if data.get('subtitle'):
        bx(slide, 1.0, 8.4, 17, 1.0, data['subtitle'], sz=22,
           color=st['chrome_muted'], font=BF)


def _section_neongrid(slide, data, st):
    F = st['title_font']; MF = st['mono_font']
    num = str(data.get('num', '01'))
    bx(slide, 1.0, 2.2, 8, 1.0, f'// SECTION {num}', sz=18, bold=True,
       color=st['accent'], font=MF)
    bx(slide, 1.0, 3.6, 18, 3.0, data.get('title', 'Section'),
       sz=72, bold=True, color=st['chrome_text'], font=F)
    rc(slide, 1.0, 7.0, 4.5, 0.06, fill=st['accent'])
    if data.get('subtitle'):
        bx(slide, 1.0, 7.4, 17, 1.4, data['subtitle'], sz=20,
           color=st['chrome_muted'], font=MF)


def _section_magazine(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    num = str(data.get('num', '01'))
    bx(slide, 1.5, 1.2, 17, 0.5, f'CHAPTER {num}', sz=14, bold=True,
       color=st['accent'], align=PP_ALIGN.CENTER, font=F)
    rc(slide, 9.0, 1.95, 2.0, 0.03, fill=st['accent'])
    bx(slide, 1.5, 4.0, 17, 3.5, data.get('title', 'Section'),
       sz=76, bold=True, italic=True,
       color=st['chrome_text'], align=PP_ALIGN.CENTER, font=F)
    if data.get('subtitle'):
        bx(slide, 2.5, 7.8, 15, 1.4, data['subtitle'], sz=20, italic=True,
           color=st['chrome_muted'], align=PP_ALIGN.CENTER, font=BF)


def _section_minimal(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    num = str(data.get('num', '01'))
    bx(slide, 1.1, 1.4, 4, 0.45, num, sz=14, bold=True,
       color=st['accent'], font=F)
    bx(slide, 1.1, 4.6, 18, 3.8, data.get('title', 'Section'),
       sz=72, color=st['chrome_text'], font=F)
    rc(slide, 1.1, 8.5, 0.6, 0.04, fill=st['accent'])
    if data.get('subtitle'):
        bx(slide, 1.1, 8.8, 15, 1.0, data['subtitle'], sz=18,
           color=st['chrome_muted'], font=BF)


def _section_brutalist(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    num = str(data.get('num', '01'))
    # Massive number block
    rc(slide, 0, 0.8, 7.5, 6.0, fill=st['accent'])
    bx(slide, 0.3, 0.9, 7.0, 6.0, num, sz=260, bold=True,
       color=st['navy'], align=PP_ALIGN.CENTER, font=F)
    bx(slide, 8.2, 3.2, 11, 4.0, (data.get('title', 'SECTION')).upper(),
       sz=68, bold=True, color=st['chrome_text'], font=F)
    if data.get('subtitle'):
        bx(slide, 8.2, 7.8, 11, 1.5, (data['subtitle']).upper(),
           sz=20, bold=True, color=st['chrome_muted'], font=BF)


_SECTION_COMPOSERS = {
    'classic':   _section_classic,
    'gradient':  _section_gradient,
    'neon-grid': _section_neongrid,
    'magazine':  _section_magazine,
    'minimal':   _section_minimal,
    'brutalist': _section_brutalist,
}


def render_section_divider(slide, data, st, _n, _total):
    # bg_color override (v1 feature) — if user sets it, fall back to solid fill
    if 'bg_color' in data:
        rc(slide, 0, 0, SW, SH, fill=_rgb(data['bg_color']))
    else:
        _paint_chrome_bg(slide, st, 'section')
    _paint_chrome_motif(slide, st, 'section')
    composer = _SECTION_COMPOSERS.get(st.get('chrome_style', 'classic'),
                                      _section_classic)
    composer(slide, data, st)


# ── Ending composers ─────────────────────────────────────────────────────────

def _ending_actions_row(slide, data, st, y_rule=7.1, divider_c=None):
    """Shared helper: draws actions strip across the bottom of an ending slide."""
    F = st['title_font']; BF = st['body_font']
    acts = data.get('actions', [])
    if not acts: return
    rc(slide, 0, y_rule, SW, 0.03, fill=_rgb('FFFFFF'))
    _set_transparency(rc(slide, 0, y_rule, SW, 0.03, fill=_rgb('FFFFFF')), 88)
    aw = SW / max(len(acts), 1)
    for i, a in enumerate(acts):
        ax = i * aw
        if i > 0:
            sep = rc(slide, ax, y_rule + 0.04, 0.03, 4.0, fill=_rgb('FFFFFF'))
            _set_transparency(sep, 88)
        bx(slide, ax + 0.7, y_rule + 0.26, 1.1, 1.0,
           a.get('icon', ''), sz=44, align=PP_ALIGN.CENTER, font=F)
        bx(slide, ax + 2.0, y_rule + 0.30, aw - 2.2, 0.7,
           a.get('title', ''), sz=24, bold=True,
           color=divider_c or st['accent'], font=F)
        bx(slide, ax + 2.0, y_rule + 1.04, aw - 2.2, 1.3,
           a.get('desc', ''), sz=21, color=st['chrome_muted'], font=BF)


def _ending_classic(slide, data, st):
    F = st['title_font']; M = st['motifs']
    if not st.get('motif_spec'):
        ov(slide, 5.0, 0.6, 10, 5.6, fill=st['accent2'], transparency=93)
    bx(slide, 1.0, 1.4, 18, 2.2, data.get('title', 'Thank You'),
       sz=72, color=st['chrome_text'], align=PP_ALIGN.CENTER, font=F)
    rc(slide, 8.9, 3.8, 2.2, 0.1,
       fill=_rgb(M.get('divider_color', '0076B3')))
    if data.get('subtitle'):
        bx(slide, 2.0, 4.1, 16, 0.9, data['subtitle'],
           sz=28, color=st['chrome_muted'],
           align=PP_ALIGN.CENTER, font=st['body_font'])
    _ending_actions_row(slide, data, st, y_rule=7.1,
                        divider_c=_rgb(M.get('number_color', '0076B3')))


def _ending_gradient(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    bx(slide, 1.0, 1.6, 18, 2.8, data.get('title', 'Thank You'),
       sz=96, bold=True, color=st['chrome_text'],
       align=PP_ALIGN.CENTER, font=F)
    rc(slide, 9.0, 4.6, 2.0, 0.06, fill=st['accent'])
    if data.get('subtitle'):
        bx(slide, 2.0, 4.9, 16, 1.4, data['subtitle'], sz=26,
           color=st['chrome_muted'], align=PP_ALIGN.CENTER, font=BF)
    _ending_actions_row(slide, data, st, y_rule=7.3)


def _ending_neongrid(slide, data, st):
    F = st['title_font']; MF = st['mono_font']
    bx(slide, 1.0, 1.2, 18, 0.5, '// END OF LINE', sz=16, bold=True,
       color=st['accent'], align=PP_ALIGN.CENTER, font=MF)
    bx(slide, 1.0, 2.2, 18, 3.2, data.get('title', 'Thank You.'),
       sz=96, bold=True, color=st['chrome_text'],
       align=PP_ALIGN.CENTER, font=F)
    rc(slide, 8.5, 5.8, 3.0, 0.06, fill=st['accent'])
    if data.get('subtitle'):
        bx(slide, 1.0, 6.1, 18, 1.0, data['subtitle'], sz=22,
           color=st['chrome_muted'], align=PP_ALIGN.CENTER, font=MF)
    _ending_actions_row(slide, data, st, y_rule=7.4)


def _ending_magazine(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    rc(slide, 1.5, 1.0, SW - 3.0, 0.02, fill=st['text_dim'])
    bx(slide, 1.5, 1.15, 18, 0.45, 'FINIS', sz=12, bold=True,
       color=st['text_dim'], align=PP_ALIGN.CENTER, font=F)
    bx(slide, 1.0, 3.0, 18, 3.5, data.get('title', 'Thank You'),
       sz=96, bold=True, italic=True, color=st['chrome_text'],
       align=PP_ALIGN.CENTER, font=F)
    rc(slide, 9.0, 6.4, 2.0, 0.03, fill=st['accent'])
    if data.get('subtitle'):
        bx(slide, 2.0, 6.7, 16, 1.2, data['subtitle'], sz=22, italic=True,
           color=st['chrome_muted'], align=PP_ALIGN.CENTER, font=BF)
    _ending_actions_row(slide, data, st, y_rule=7.5)


def _ending_minimal(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    bx(slide, 1.1, 4.0, 18, 2.4, data.get('title', 'Thank You.'),
       sz=72, color=st['chrome_text'], font=F)
    rc(slide, 1.1, 6.5, 0.6, 0.04, fill=st['accent'])
    if data.get('subtitle'):
        bx(slide, 1.1, 6.8, 15, 1.0, data['subtitle'], sz=20,
           color=st['chrome_muted'], font=BF)
    _ending_actions_row(slide, data, st, y_rule=8.4)


def _ending_brutalist(slide, data, st):
    F = st['title_font']; BF = st['body_font']
    rc(slide, 0, 2.0, SW, 3.6, fill=st['accent'])
    bx(slide, 1.0, 2.1, 18, 3.3, (data.get('title', 'THANK YOU')).upper(),
       sz=120, bold=True, color=st['navy'],
       align=PP_ALIGN.CENTER, font=F)
    if data.get('subtitle'):
        bx(slide, 1.0, 5.9, 18, 1.2, (data['subtitle']).upper(),
           sz=24, bold=True, color=st['chrome_text'],
           align=PP_ALIGN.CENTER, font=BF)
    _ending_actions_row(slide, data, st, y_rule=7.6)


_ENDING_COMPOSERS = {
    'classic':   _ending_classic,
    'gradient':  _ending_gradient,
    'neon-grid': _ending_neongrid,
    'magazine':  _ending_magazine,
    'minimal':   _ending_minimal,
    'brutalist': _ending_brutalist,
}


def render_ending(slide, data, st, n, total):
    _paint_chrome_bg(slide, st, 'ending')
    _paint_chrome_motif(slide, st, 'ending')
    composer = _ENDING_COMPOSERS.get(st.get('chrome_style', 'classic'),
                                     _ending_classic)
    composer(slide, data, st)
    _slide_num_free(slide, n, total, st=st)


# ─────────────────────────────────────────────────────────────────────────────
# Content slide renderers
# ─────────────────────────────────────────────────────────────────────────────

def render_bullets_with_panel(slide, data, st, num, total):
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    bullets = data.get('bullets', [])
    panel = data.get('side_panel', {})
    has_panel = panel.get('type') == 'quote' and panel.get('text')
    bullet_w = 9.0 if has_panel else 17.0  # full width when no panel

    n_bullets = len(bullets)
    avail_h = SH - ct - 0.4
    bullet_spacing = avail_h / max(n_bullets, 1)
    bullet_spacing = min(bullet_spacing, 2.0)
    block_h = n_bullets * bullet_spacing
    by = ct + (avail_h - block_h) / 2 + 0.1
    for b in bullets:
        rc(slide, 1.0, by + 0.12, 0.14, 0.14, fill=st['accent'])
        bx(slide, 1.32, by, bullet_w, 0.8, b, sz=20,
           color=st['text_white'] if blue else st['text_dark'], font=F)
        by += bullet_spacing

    if has_panel:
        panel_h = SH - ct - 0.4
        rc(slide, 11.0, ct + 0.1, 7.9, panel_h, fill=st['card_side'])
        rc(slide, 11.0, ct + 0.1, 0.12, panel_h, fill=st['accent'])
        bx(slide, 11.35, ct + 0.4, 7.3, 1.1, '\u201C', sz=80, italic=True,
           color=_rgb([0x33, 0x88, 0xBB]), font=F)
        bx(slide, 11.35, ct + 1.5, 7.3, panel_h - 1.8,
           panel.get('text', ''), sz=24, italic=True,
           color=st['white'], font=F)

    rules = data.get('ground_rules', [])
    if rules:
        rc(slide, 1.0, by + 0.12, 9.5, 2.55, fill=st['card'])
        bx(slide, 1.15, by + 0.20, 4, 0.28, 'GROUND RULES', sz=8,
           bold=True, color=st['accent'], font=F)
        ry = by + 0.58
        for rt in rules:
            bx(slide, 1.18, ry, 0.28, 0.30, '—', sz=13, bold=True,
               color=st['accent'], font=F)
            bx(slide, 1.50, ry, 8.8, 0.40, rt, sz=13,
               color=st['white'], font=F)
            ry += 0.48


def render_two_column_check(slide, data, st, num, total):
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct - 0.05, 17, 0.35, data['subtitle'], sz=15,
           color=st['text_gray'], font=F)

    LIGHT_COLORS = {
        'green': [0xEB, 0xF9, 0xF0], 'red': [0xFF, 0xEB, 0xEB],
        'blue':  [0xEF, 0xF7, 0xFD], 'yellow': [0xFF, 0xFB, 0xE6],
    }
    col_h = SH - ct - 0.3
    for col_i, side in enumerate([data.get('left', {}), data.get('right', {})]):
        x = 1.0 + col_i * 9.5
        hc = _rgb(side.get('color', [0x27, 0xAE, 0x60] if col_i == 0
                  else [0xE7, 0x4C, 0x3C]))
        lc_key = side.get('light_color', 'green' if col_i == 0 else 'red')
        lc = _rgb(LIGHT_COLORS.get(lc_key, LIGHT_COLORS['green']))
        marker = side.get('marker', '✓' if col_i == 0 else '✗')
        rc(slide, x, ct + 0.3, 8.8, col_h, fill=lc)
        rc(slide, x, ct + 0.3, 8.8, 0.58, fill=hc)
        bx(slide, x + 0.20, ct + 0.34, 8.3, 0.50, side.get('title', ''),
           sz=16, bold=True, color=st['white'], font=F)
        items = side.get('items', [])
        n_items = len(items)
        item_space = (col_h - 0.85) / max(n_items, 1)
        iy = ct + 1.08
        for item in items:
            bx(slide, x + 0.15, iy, 0.40, 0.40, marker, sz=19, bold=True,
               color=hc, font=F)
            bx(slide, x + 0.60, iy, 8.0, 0.70, item, sz=19,
               color=st['text_dark'] if not blue else st['white'], font=F)
            iy += item_space


def render_cards_grid(slide, data, st, num, total):
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    cards = data.get('cards', [])
    n = len(cards)
    if n == 0: return
    CW = 8.7; GX = 0.5; GY = 0.28
    rows_n = math.ceil(n / 2)
    avail_h = SH - ct - 0.2
    CH = min(2.55, max(1.4, (avail_h - GY * (rows_n - 1)) / rows_n))

    for i, card in enumerate(cards[:8]):
        col = i % 2; row = i // 2
        x = 1.0 + col * (CW + GX)
        y = ct + 0.1 + row * (CH + GY)
        bg = st['card'] if blue else st['card_white']
        rc(slide, x, y, CW, CH, fill=bg)
        bx(slide, x + 0.2, y + CH * 0.06, 1.5, CH * 0.38, card.get('num', ''),
           sz=32, bold=True, color=st['card_num'], font=F)
        bx(slide, x + 0.2, y + CH * 0.48, CW - 0.4, CH * 0.48,
           card.get('text', ''), sz=18,
           color=st['text_white'] if blue else st['text_dark'], font=F)


def render_criteria_rows(slide, data, st, num, total):
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct - 0.05, 17, 0.35, data['subtitle'], sz=15,
           color=st['text_dim'] if blue else st['text_gray'], font=F)

    rows = data.get('criteria', [])
    n = len(rows)
    if n == 0: return
    gap = 0.25
    avail_h = SH - ct - 0.3
    RH = min(2.6, max(1.2, (avail_h - gap * (n - 1)) / n))
    block_h = n * RH + (n - 1) * gap
    start_y = ct + (avail_h - block_h) / 2 + 0.1

    for i, row in enumerate(rows):
        y = start_y + i * (RH + gap)
        bg = st['card'] if blue else st['card_white']
        rc(slide, 1.0, y, 17.5, RH, fill=bg)
        bx(slide, 1.1, y + RH * 0.15, 1.5, RH * 0.54, row.get('num', ''),
           sz=36, bold=True, color=st['card_num'], font=F)
        rc(slide, 2.7, y + RH * 0.12, 0.04, RH * 0.76, fill=st['line_gray'])
        bx(slide, 2.9, y + RH * 0.08, 5.5, RH * 0.32, row.get('label', ''),
           sz=14, bold=True, color=st['accent'], font=F)
        bx(slide, 2.9, y + RH * 0.40, 15.0, RH * 0.55, row.get('text', ''),
           sz=20, color=st['text_white'] if blue else st['text_dark'], font=F)

    if data.get('footnote'):
        fn_y = ct + 0.4 + n * (RH + gap) - 0.15
        bx(slide, 1.0, fn_y, 17.5, 0.35, data['footnote'], sz=13,
           color=st['text_dim'] if blue else st['text_gray'],
           align=PP_ALIGN.RIGHT, font=F)


def render_scope_tiers(slide, data, st, num, total):
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct - 0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_dim'] if blue else st['text_gray'], font=F)

    tiers = data.get('tiers', [])
    n = len(tiers)
    if n == 0: return
    gap = 0.18
    avail_h = SH - ct - 0.45
    BH = min(2.2, max(1.0, (avail_h - gap * (n - 1)) / n))

    for i, tier in enumerate(tiers):
        y = ct + 0.35 + i * (BH + gap)
        tc = _rgb(tier.get('color', [0x44, 0x99, 0x44]))
        bg = st['card'] if blue else st['card_white']
        rc(slide, 1.0, y, 17.5, BH, fill=bg)
        rc(slide, 1.0, y, 0.14, BH, fill=tc)
        bx(slide, 1.28, y + BH * 0.08, 1.0, BH * 0.43, tier.get('icon', ''),
           sz=22, color=st['white'], font=F)
        bx(slide, 2.45, y + BH * 0.06, 15.3, BH * 0.20, tier.get('label', ''),
           sz=9, bold=True, color=st['accent'], font=F)
        bx(slide, 2.45, y + BH * 0.28, 15.3, BH * 0.26, tier.get('desc', ''),
           sz=13.5, color=st['text_white'] if blue else st['text_dark'], font=F)
        bx(slide, 2.45, y + BH * 0.58, 15.3, BH * 0.30, tier.get('example', ''),
           sz=11, italic=True,
           color=st['text_dim'] if blue else st['text_gray'], font=F)


def render_two_panel(slide, data, st, num, total):
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    panel_h = SH - ct - 0.3
    for col_i, panel in enumerate(data.get('panels', [])[:2]):
        x = 1.0 + col_i * 9.5
        hc = _rgb(panel.get('color', [0x00, 0x76, 0xB3]))
        bg = st['card'] if blue else st['card_white']
        rc(slide, x, ct + 0.15, 8.8, panel_h, fill=bg)
        rc(slide, x, ct + 0.15, 8.8, 0.60, fill=st['navy'])
        bx(slide, x + 0.20, ct + 0.22, 8.3, 0.50,
           f"{panel.get('icon', '')}  {panel.get('title', '')}", sz=16,
           bold=True, color=st['white'], font=F)
        items = panel.get('items', [])
        n_items = len(items)
        content_area = panel_h - 0.85
        item_space = content_area / max(n_items, 1)
        item_space = min(item_space, 2.0)  # cap spacing
        block_h = n_items * item_space
        iy = ct + 0.95 + (content_area - block_h) / 2
        for item in items:
            bx(slide, x + 0.15, iy, 0.40, 0.40, '▸', sz=19, bold=True,
               color=hc, font=F)
            bx(slide, x + 0.60, iy, 8.0, 0.85, item, sz=19,
               color=st['text_white'] if blue else st['text_dark'], font=F)
            iy += item_space


def render_two_column_steps(slide, data, st, num, total):
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    steps_list = data.get('columns', [])[:2]
    # Calculate card height to fill available space
    max_steps = max((len(c.get('steps', [])) for c in steps_list), default=3)
    avail_h = SH - ct - 0.9
    card_gap = 0.22
    card_h = min(2.6, max(1.6, (avail_h - 0.5 - card_gap * (max_steps - 1)) / max_steps))

    for col_i, col in enumerate(steps_list):
        x = 1.0 + col_i * 9.5
        bx(slide, x, ct + 0.08, 9, 0.40, col.get('title', ''), sz=15,
           bold=True, color=st['accent'], font=F)
        sy = ct + 0.60
        for n, step in enumerate(col.get('steps', [])[:4], 1):
            bg = st['card'] if blue else st['card_white']
            rc(slide, x, sy, 8.8, card_h, fill=bg)
            rc(slide, x, sy, 0.55, card_h, fill=st['card_side'])
            bx(slide, x + 0.08, sy + card_h * 0.32, 0.40, 0.55, str(n), sz=17,
               bold=True, color=st['white'], align=PP_ALIGN.CENTER, font=F)
            tb = bx(slide, x + 0.72, sy + card_h * 0.12, 7.9, card_h * 0.76,
                    step.get('bold', ''), sz=19, bold=True,
                    color=st['text_white'] if blue else st['text_dark'], font=F)
            ap(tb.text_frame, step.get('normal', ''), sz=17,
               color=st['text_dim'] if blue else st['text_gray'], font=F)
            sy += card_h + card_gap

    if data.get('warning'):
        rc(slide, 1.0, 8.60, 17.5, 0.72, fill=_rgb([0x3A, 0x08, 0x08]))
        rc(slide, 1.0, 8.60, 0.14, 0.72, fill=_rgb([0xE7, 0x4C, 0x3C]))
        bx(slide, 1.25, 8.66, 17.0, 0.55, data['warning'], sz=12,
           color=_rgb([0xFF, 0xAA, 0xAA]), font=F)


def render_scenario_cards(slide, data, st, num, total):
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct - 0.05, 17, 0.35, data['subtitle'], sz=15,
           color=st['text_dim'] if blue else st['text_gray'], font=F)

    SW_C = 8.7; SH_C = 4.0
    for i, scen in enumerate(data.get('scenarios', [])[:4]):
        col = i % 2; row = i // 2
        x = 1.0 + col * (SW_C + 0.5)
        y = ct + 0.25 + row * (SH_C + 0.25)
        bg = st['card'] if blue else st['card_white']
        rc(slide, x, y, SW_C, SH_C, fill=bg)
        rc(slide, x, y, SW_C, 0.09, fill=st['accent'])
        bx(slide, x + 0.25, y + 0.20, 3.0, 0.35,
           f"Scenario {scen.get('num', '')}", sz=12, bold=True,
           color=st['accent'], font=F)
        bx(slide, x + 0.25, y + 0.58, SW_C - 0.5, 0.65,
           f"{scen.get('icon', '')}  {scen.get('title', '')}", sz=20,
           bold=True, color=st['text_white'] if blue else st['text_dark'], font=F)
        bx(slide, x + 0.25, y + 1.35, SW_C - 0.5, 2.0,
           scen.get('desc', ''), sz=18,
           color=st['text_dim'] if blue else st['text_gray'], font=F)
        rc(slide, x + 0.25, y + SH_C - 0.55, 2.5, 0.38, fill=st['card_side'])
        bx(slide, x + 0.31, y + SH_C - 0.51, 2.35, 0.32,
           scen.get('tag', ''), sz=12, color=st['accent'], font=F)


def render_survey(slide, data, st, num, total):
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct - 0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_gray'], font=F)

    STH = 1.78
    for i, step in enumerate(data.get('steps', [])[:4]):
        y = ct + 0.35 + i * (STH + 0.22)
        bg = st['card_white'] if not blue else st['card']
        rc(slide, 1.0, y, 10.5, STH, fill=bg)
        rc(slide, 1.0, y, 0.65, STH, fill=st['accent2'])
        bx(slide, 1.06, y + 0.56, 0.53, 0.65, str(i + 1), sz=22,
           bold=True, color=st['white'], align=PP_ALIGN.CENTER, font=F)
        bx(slide, 1.82, y + 0.26, 9.5, 0.48, step.get('title', ''),
           sz=15, bold=True,
           color=st['text_dark'] if not blue else st['white'], font=F)
        bx(slide, 1.82, y + 0.84, 9.5, 0.72, step.get('desc', ''),
           sz=13, color=st['text_gray'], font=F)

    qr_bg = st['card_white'] if not blue else st['card']
    rc(slide, 12.3, ct + 0.3, 6.6, 9.0, fill=qr_bg,
       line=st['accent2'], lw=1.0)
    bx(slide, 12.3, ct + 0.45, 6.6, 0.34,
       data.get('qr_label', 'SURVEY QR CODE'), sz=9, bold=True,
       color=st['accent2'], align=PP_ALIGN.CENTER, font=F)
    rc(slide, 13.25, ct + 1.05, 4.6, 4.8,
       fill=_rgb([0xD8, 0xEA, 0xF5]), line=st['accent2'], lw=0.75)
    bx(slide, 13.25, ct + 2.95, 4.6, 1.0, '[ QR Code ]', sz=14,
       italic=True, color=_rgb([0x66, 0x99, 0xBB]),
       align=PP_ALIGN.CENTER, font=F)
    bx(slide, 12.3, ct + 6.1, 6.6, 0.50,
       data.get('qr_note', '(Link provided on the day)'), sz=11,
       italic=True, color=_rgb([0x77, 0x99, 0xBB]),
       align=PP_ALIGN.CENTER, font=F)
    bx(slide, 12.3, ct + 6.8, 6.6, 0.80,
       data.get('qr_hint',
                'Moderator will circulate\nto assist with submission'),
       sz=12, color=st['text_gray'], align=PP_ALIGN.CENTER, font=F)


# ── NEW: stat_highlight ───────────────────────────────────────────────────────

def render_stat_highlight(slide, data, st, num, total):
    """2–4 large statistics displayed as bold data cards in a row."""
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    stats = data.get('stats', [])
    n = len(stats)
    if n == 0: return

    gap   = 0.4
    CW    = (SW - 2.0 - gap * (n - 1)) / n
    CH    = SH - ct - 0.5
    y     = ct + 0.25

    for i, stat in enumerate(stats):
        x  = 1.0 + i * (CW + gap)
        bg = st['card'] if blue else st['card_white']
        rc(slide, x, y, CW, CH, fill=bg)
        rc(slide, x, y, CW, 0.12, fill=st['accent'])

        bx(slide, x + 0.15, y + CH * 0.12, CW - 0.3, CH * 0.35,
           stat.get('value', ''), sz=72, bold=True,
           color=st['accent'], align=PP_ALIGN.CENTER, font=F)

        bx(slide, x + 0.15, y + CH * 0.50, CW - 0.3, CH * 0.18,
           stat.get('label', ''), sz=24, bold=True,
           color=st['text_white'] if blue else st['text_dark'],
           align=PP_ALIGN.CENTER, font=F)

        bx(slide, x + 0.15, y + CH * 0.68, CW - 0.3, CH * 0.20,
           stat.get('desc', ''), sz=17,
           color=st['text_dim'] if blue else st['text_gray'],
           align=PP_ALIGN.CENTER, font=F)

        # Trend arrow (optional)
        trend = stat.get('trend', '')
        if trend:
            arrow_map = {'up': ('▲', '27AE60'), 'down': ('▼', 'E74C3C'),
                         'flat': ('—', '999999')}
            sym, clr = arrow_map.get(trend, ('', '999999'))
            if sym:
                bx(slide, x + 0.15, y + CH * 0.88, CW - 0.3, CH * 0.10,
                   sym, sz=18, bold=True, color=_rgb(clr),
                   align=PP_ALIGN.CENTER, font=F)


# ── NEW: timeline ─────────────────────────────────────────────────────────────

def render_timeline(slide, data, st, num, total):
    """Horizontal timeline with cards: dot on line, card below with content."""
    blue   = data.get('background', 'blue') == 'blue'
    F      = st['font']
    M      = st['motifs']
    ct     = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
                 data.get('title', ''), blue=blue, st=st)

    items  = data.get('items', [])
    n      = len(items)
    if n == 0: return

    lc     = st['accent']
    nc     = _rgb(M.get('number_color', st['accent']))
    tc     = st['text_white'] if blue else st['text_dark']
    dc     = st['text_dim'] if blue else st['text_gray']
    card_bg = st['card'] if blue else st['card_white']

    # Layout — center the whole assembly vertically
    dot_r    = 0.28
    gap      = 0.35
    slot_w   = 18.0 / n
    card_w   = slot_w - gap
    card_h   = 3.8                # fixed height, enough for content
    assembly_h = dot_r * 2 + 0.30 + card_h  # dot + gap + card
    avail_h  = SH - ct - 0.3
    top_pad  = (avail_h - assembly_h) / 2.0
    line_y   = ct + top_pad + dot_r
    card_top = line_y + dot_r + 0.30

    # Horizontal line
    rc(slide, 1.0, line_y - 0.045, 18.0, 0.09, fill=lc)

    for i, item in enumerate(items):
        cx = 1.0 + slot_w * i + slot_w / 2.0
        card_x = cx - card_w / 2.0

        # Dot on line
        s = slide.shapes.add_shape(9,
            I(cx - dot_r), I(line_y - dot_r),
            I(dot_r * 2),  I(dot_r * 2))
        s.fill.solid(); s.fill.fore_color.rgb = lc
        s.line.fill.background()

        # Card below
        rc(slide, card_x, card_top, card_w, card_h, fill=card_bg)
        rc(slide, card_x, card_top, card_w, 0.08, fill=lc)

        # Phase label (inside card top)
        bx(slide, card_x + 0.2, card_top + 0.25, card_w - 0.4, 0.50,
           item.get('phase', ''), sz=14, bold=True,
           color=nc, align=PP_ALIGN.CENTER, font=F)

        # Title (large, centered in card)
        bx(slide, card_x + 0.2, card_top + 0.85, card_w - 0.4, 1.0,
           item.get('title', ''), sz=24, bold=True,
           color=tc, align=PP_ALIGN.CENTER, font=F)

        # Accent divider inside card
        div_w = 0.8
        rc(slide, cx - div_w / 2, card_top + 1.95, div_w, 0.05, fill=lc)

        # Description
        bx(slide, card_x + 0.2, card_top + 2.15, card_w - 0.4, card_h - 2.4,
           item.get('desc', ''), sz=17,
           color=dc, align=PP_ALIGN.CENTER, font=F)


# ── NEW: quote_full ───────────────────────────────────────────────────────────

def render_quote_full(slide, data, st, num, total):
    """Full-slide featured quote with large decorative treatment."""
    blue = data.get('background', 'blue') == 'blue'
    F    = st['font']

    # Minimal header (no rule line — quote needs clean space)
    lc = st['accent'] if blue else st['accent2']
    nc = st['text_num'] if blue else _rgb('AAAAAA')
    bx(slide, 1.0, 0.35, 10, 0.28, data.get('section', ''), sz=8, color=lc, font=F)
    bx(slide, 16.5, 0.35, 2.5, 0.28,
       _slide_num(data, num, total), sz=8,
       color=nc, align=PP_ALIGN.RIGHT, font=F)

    # Left accent bar
    rc(slide, 0.8, 1.8, 0.14, 7.2, fill=st['accent'])

    # Large decorative opening quote mark
    bx(slide, 1.5, 1.0, 4.5, 2.8, '\u201C', sz=120, italic=True,
       color=st['accent2'], font=F)

    # Quote body
    bx(slide, 1.5, 3.0, 16.5, 5.2,
       data.get('quote', ''), sz=28, italic=True,
       color=st['text_white'] if blue else st['text_dark'],
       align=PP_ALIGN.LEFT, font=F)

    # Attribution
    if data.get('attribution'):
        rc(slide, 1.5, 8.5, 16.5, 0.04, fill=st['line_gray'])
        bx(slide, 1.5, 8.7, 16.5, 0.65,
           data['attribution'], sz=14,
           color=st['text_dim'] if blue else st['text_gray'],
           align=PP_ALIGN.LEFT, font=F)

    _slide_num_free(slide, num, total, light=blue, st=st)


# ── NEW: center_focus ─────────────────────────────────────────────────────────

def render_center_focus(slide, data, st, num, total):
    """Center-dominant layout for key statements, strategic focus, section openers."""
    blue = data.get('background', 'blue') == 'blue'
    F    = st['font']
    M    = st['motifs']

    # Small section label top-left
    lc = st['accent'] if blue else st['accent2']
    bx(slide, 1.0, 0.35, 12, 0.30, data.get('section', ''), sz=8, color=lc, font=F)

    # Large decorative background circle
    ov(slide, 3.5, 0.8, 13.0, 9.5, fill=st['accent2'], transparency=91)

    # Optional title label (small caps, above message)
    if data.get('title'):
        bx(slide, 1.5, 2.6, 17.0, 0.48,
           data['title'].upper(), sz=10, bold=True,
           color=st['accent'], align=PP_ALIGN.CENTER, font=F)

    # Main message — large, centered
    bx(slide, 1.5, 3.3, 17.0, 3.8,
       data.get('message', ''), sz=36, bold=True,
       color=st['text_white'] if blue else st['text_dark'],
       align=PP_ALIGN.CENTER, font=F)

    # Accent divider
    rc(slide, 8.75, 7.3, 2.5, 0.1,
       fill=_rgb(M.get('divider_color', st['accent'])))

    # Optional context line
    if data.get('context'):
        bx(slide, 1.5, 7.6, 17.0, 0.75,
           data['context'], sz=16,
           color=st['text_dim'] if blue else st['text_gray'],
           align=PP_ALIGN.CENTER, font=F)

    _slide_num_free(slide, num, total, light=blue, st=st)


# ── NEW: comparison_table ─────────────────────────────────────────────────────

def render_comparison_table(slide, data, st, num, total):
    """Structured comparison table using native PowerPoint table object."""
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    from lxml import etree

    blue = data.get('background', 'white') == 'blue'
    F    = st['font']
    ct   = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
               data.get('title', ''), blue=blue, st=st)

    cols = data.get('columns', [])
    rows = data.get('rows', [])
    if not cols or not rows: return

    nc     = len(cols)
    n_rows = len(rows)
    n_data_rows = n_rows + 1  # +1 for header

    # Table geometry
    tbl_left = 1.0
    tbl_top  = ct + 0.3
    tbl_w    = SW - 2.0  # 18" usable
    avail_h  = SH - tbl_top - 0.3
    row_h    = min(1.55, max(0.65, avail_h / n_data_rows))
    tbl_h    = row_h * n_data_rows

    # Column widths: label column + N data columns
    label_w = 3.5
    data_w  = (tbl_w - label_w) / nc

    # Create native table
    tbl_shape = slide.shapes.add_table(
        n_data_rows, nc + 1,
        I(tbl_left), I(tbl_top), I(tbl_w), I(tbl_h))
    table = tbl_shape.table

    # Set column widths
    table.columns[0].width = I(label_w)
    for ci in range(nc):
        table.columns[ci + 1].width = I(data_w)

    # Set row heights
    for ri in range(n_data_rows):
        table.rows[ri].height = I(row_h)

    # Helper: set cell fill
    def _cell_fill(cell, color):
        tcPr = cell._tc.get_or_add_tcPr()
        sf = etree.SubElement(tcPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        r, g, b = color
        sc.set('val', f'{r:02X}{g:02X}{b:02X}')

    # Helper: set cell text
    def _cell_text(cell, text, bold=False, sz=19, align=PP_ALIGN.LEFT,
                   color=None):
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = F
        run.font.size = P(sz)
        run.font.bold = bold
        if color:
            run.font.color.rgb = _rgb(color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left  = I(0.15)
        cell.margin_right = I(0.15)
        cell.margin_top   = I(0.08)
        cell.margin_bottom= I(0.08)

    # Helper: set thin borders on a cell
    def _cell_border(cell, color_hex='3A3A45'):
        tcPr = cell._tc.get_or_add_tcPr()
        for tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = etree.SubElement(tcPr, qn(tag))
            ln.set('w', str(int(P(0.5))))
            sf = etree.SubElement(ln, qn('a:solidFill'))
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', color_hex)

    # Determine text/bg colors
    txt_hdr  = st['white']
    txt_body = st['text_white'] if blue else st['text_dark']
    border_c = '%02X%02X%02X' % (st['line_gray'][0], st['line_gray'][1],
                                  st['line_gray'][2]) if hasattr(st['line_gray'], '__getitem__') else '3A3A45'

    # Fill header row
    hdr_cell = table.cell(0, 0)
    _cell_fill(hdr_cell, st['accent'])
    _cell_text(hdr_cell, '', bold=True, color=txt_hdr)
    _cell_border(hdr_cell, border_c)

    for ci, col_title in enumerate(cols):
        cell = table.cell(0, ci + 1)
        hdr_bg = st['accent2'] if ci % 2 == 0 else st['card_side']
        _cell_fill(cell, hdr_bg)
        _cell_text(cell, col_title, bold=True, sz=20,
                   align=PP_ALIGN.CENTER, color=txt_hdr)
        _cell_border(cell, border_c)

    # Fill data rows
    for ri, row in enumerate(rows):
        if blue:
            row_bg = st['card'] if ri % 2 == 0 else st['card_side']
        else:
            row_bg = _rgb('EEF4FC') if ri % 2 == 0 else st['card_white']

        # Label cell
        cell = table.cell(ri + 1, 0)
        _cell_fill(cell, row_bg)
        _cell_text(cell, row.get('label', ''), bold=True, sz=19,
                   color=txt_body)
        _cell_border(cell, border_c)

        # Value cells
        for ci, val in enumerate(row.get('values', [])[:nc]):
            cell = table.cell(ri + 1, ci + 1)
            _cell_fill(cell, row_bg)
            _cell_text(cell, val, sz=19, align=PP_ALIGN.CENTER,
                       color=txt_body)
            _cell_border(cell, border_c)


# ─────────────────────────────────────────────────────────────────────────────
# Validate and auto-fix overflow
# ─────────────────────────────────────────────────────────────────────────────

def _estimate_overflow(tf, box_w_in, box_h_in):
    """Estimate if text frame overflows its bounding box. Returns True if likely."""
    total_h = 0.0
    for para in tf.paragraphs:
        if not para.runs:
            total_h += 0.18
            continue
        pt = para.runs[0].font.size.pt if para.runs[0].font.size else 13.0
        chars_per_line = max(1, int(box_w_in * 96.0 / (pt * 0.58)))
        text_len = len(para.text)
        lines = max(1, math.ceil(text_len / chars_per_line)) if text_len else 1
        total_h += lines * pt * 1.35 / 72.0
    return total_h > box_h_in * 0.90

def _shrink_font(tf, delta_pt=1):
    """Reduce all run font sizes in a text frame by delta_pt (floor: 8pt)."""
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size:
                run.font.size = P(max(8, run.font.size.pt - delta_pt))

def validate_and_fix(prs):
    """
    Scan every text box for likely overflow and shrink font to fit.
    Applies up to 4 pt reduction per shape. Logs adjustments.
    """
    fixes = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # Skip shapes with no real text content
            if not any(para.text.strip() for para in tf.paragraphs):
                continue
            bw = shape.width  / 914400.0
            bh = shape.height / 914400.0
            if bw < 0.5 or bh < 0.1:   # skip decorative micro-boxes
                continue
            for _ in range(4):
                if not _estimate_overflow(tf, bw, bh):
                    break
                _shrink_font(tf, 1)
                fixes.append(f"  slide {si + 1} '{shape.name}': -1pt")

    if fixes:
        print(f"validate_and_fix: {len(fixes)} adjustment(s) applied")
        for f in fixes:
            print(f)
    else:
        print("validate_and_fix: no overflow detected \u2713")


# ── NEW: chart ────────────────────────────────────────────────────────────────

_CHART_TYPE_MAP = {
    'bar':            XL_CHART_TYPE.COLUMN_CLUSTERED,
    'bar_stacked':    XL_CHART_TYPE.COLUMN_STACKED,
    'bar_horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
    'line':           XL_CHART_TYPE.LINE,
    'line_marker':    XL_CHART_TYPE.LINE_MARKERS,
    'area':           XL_CHART_TYPE.AREA,
    'pie':            XL_CHART_TYPE.PIE,
    'doughnut':       XL_CHART_TYPE.DOUGHNUT,
    'radar':          XL_CHART_TYPE.RADAR_FILLED,
    'scatter':        XL_CHART_TYPE.XY_SCATTER,
}

def render_chart(slide, data, st, num, total):
    """Data chart using native python-pptx charts. Supports 10 chart types."""
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section', ''), _slide_num(data, num, total),
             data.get('title', ''), blue=blue, st=st)

    chart_type_str = data.get('chart_type', 'bar')
    xl_type = _CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    series_list = data.get('series', [])
    if not series_list:
        return

    is_pie = chart_type_str in ('pie', 'doughnut')
    is_xy  = chart_type_str in ('scatter',)

    # Build chart data — scatter uses XyChartData, others use CategoryChartData
    if is_xy:
        from pptx.chart.data import XyChartData
        chart_data = XyChartData()
        for s in series_list:
            xy_series = chart_data.add_series(s.get('name', ''))
            for x, y in zip(s.get('x_values', []), s.get('y_values', [])):
                xy_series.add_data_point(x, y)
    else:
        categories = data.get('categories', [])
        if not categories:
            return
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in series_list:
            chart_data.add_series(s.get('name', ''), s.get('values', []))

    chart_left = I(1.5)
    chart_top  = I(ct + 0.3)
    chart_w    = I(17.0)
    chart_h    = I(SH - ct - 1.2)

    graphic = slide.shapes.add_chart(xl_type, chart_left, chart_top,
                                     chart_w, chart_h, chart_data)
    chart = graphic.chart

    # Apply style colors to chart series / points
    colors = st.get('chart_colors', [st['accent'], st['accent2'], st['navy']])
    if is_pie:
        # Color individual points on the first (only) series
        plot = chart.plots[0]
        for i, pt in enumerate(plot.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
    elif chart_type_str in ('line', 'line_marker'):
        for i, s in enumerate(chart.series):
            s.format.line.color.rgb = colors[i % len(colors)]
            s.format.line.width = P(2.5)
    elif chart_type_str == 'radar':
        for i, s in enumerate(chart.series):
            s.format.line.color.rgb = colors[i % len(colors)]
            s.format.line.width = P(2.0)
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[i % len(colors)]
            _set_transparency(s, 70)
    elif is_xy:
        for i, s in enumerate(chart.series):
            s.format.line.color.rgb = colors[i % len(colors)]
            s.marker.style = 8  # circle
            s.marker.size = 10
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = colors[i % len(colors)]
    else:
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = colors[i % len(colors)]

    # Chart styling: font, legend, gridlines
    chart.font.name = F
    chart.font.size = P(14)
    if not is_pie:
        chart.font.color.rgb = st['text_white'] if blue else st['text_dark']
    chart.has_legend = len(series_list) > 1 or is_pie
    if chart.has_legend:
        chart.legend.font.name = F
        chart.legend.font.size = P(14)
        chart.legend.include_in_layout = False

    # Category / value axis styling (not available for pie/radar)
    has_axes = not is_pie and chart_type_str != 'radar'
    if has_axes and not is_xy:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.name = F
        cat_ax.tick_labels.font.size = P(14)
        cat_ax.tick_labels.font.color.rgb = st['text_white'] if blue else st['text_dark']
        cat_ax.has_major_gridlines = False

        val_ax = chart.value_axis
        val_ax.tick_labels.font.name = F
        val_ax.tick_labels.font.size = P(14)
        val_ax.tick_labels.font.color.rgb = st['text_white'] if blue else st['text_dark']
        val_ax.major_gridlines.format.line.color.rgb = _rgb('3A5570') if blue else _rgb('CCCCCC')
        val_ax.major_gridlines.format.line.width = P(0.5)

    # Scatter: style both X and Y value axes
    if is_xy:
        for ax in (chart.value_axis,):
            ax.tick_labels.font.name = F
            ax.tick_labels.font.size = P(14)
            ax.tick_labels.font.color.rgb = st['text_white'] if blue else st['text_dark']

    # Optional footnote
    footnote = data.get('footnote', '')
    if footnote:
        bx(slide, 1.5, SH - 0.65, 17.0, 0.45, footnote, sz=12, italic=True,
           color=st['text_dim'] if blue else st['text_gray'],
           align=PP_ALIGN.RIGHT, font=F)


# ── NEW: image_full ──────────────────────────────────────────────────────────

def render_image_full(slide, data, st, num, total):
    """Full-bleed image with semi-transparent overlay and centered text."""
    F = st['font']
    img_path = data.get('image_path', '')
    overlay = data.get('overlay', 'dark')

    # Background image or placeholder
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, I(0), I(0), I(SW), I(SH))
    else:
        # Gradient-like placeholder
        rc(slide, 0, 0, SW, SH, fill=st['navy'])
        if img_path:
            bx(slide, 2, SH - 1.0, SW - 4, 0.6,
               f'(Image not found: {img_path})', sz=14, italic=True,
               color=_rgb('FF6666'), align=PP_ALIGN.CENTER, font=F)

    # Semi-transparent overlay
    if overlay == 'dark':
        ovl = rc(slide, 0, 0, SW, SH, fill=st['navy'])
        _set_transparency(ovl, 60)  # 40% opaque
    else:
        ovl = rc(slide, 0, 0, SW, SH, fill=_rgb('FFFFFF'))
        _set_transparency(ovl, 50)  # 50% opaque

    # Centered text
    title = data.get('title', '')
    subtitle = data.get('subtitle', '')
    txt_color = st['white'] if overlay == 'dark' else st['text_dark']

    if title:
        bx(slide, 2.0, SH * 0.32, SW - 4.0, 2.0, title, sz=36, bold=True,
           color=txt_color, align=PP_ALIGN.CENTER, font=F)
    if subtitle:
        bx(slide, 2.0, SH * 0.52, SW - 4.0, 1.5, subtitle, sz=20,
           color=txt_color, align=PP_ALIGN.CENTER, font=F)


# ── NEW: freeform ─────────────────────────────────────────────────────────────

def _resolve_color(val, st):
    """Resolve a color value: style key name, hex string, or [r,g,b] list."""
    if isinstance(val, str) and val in st:
        return st[val]
    return _rgb(val)

def render_freeform(slide, data, st, num, total):
    """Freeform slide: AI specifies exact element positions for creative layouts."""
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']

    # Optional standard header (skip if no title)
    if data.get('title'):
        hdr(slide, data.get('section', ''), _slide_num(data, num, total),
            data.get('title', ''), blue=blue, st=st)

    for el in data.get('elements', []):
        kind = el.get('kind', '')
        x, y = el.get('x', 0), el.get('y', 0)
        w, h = el.get('w', 1), el.get('h', 1)

        if kind == 'text':
            color = _resolve_color(el.get('color', 'white' if blue else 'text_dark'), st)
            align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                         'right': PP_ALIGN.RIGHT}
            bx(slide, x, y, w, h, el.get('text', ''),
               sz=el.get('sz', 18), bold=el.get('bold', False),
               italic=el.get('italic', False),
               color=color,
               align=align_map.get(el.get('align', 'left'), PP_ALIGN.LEFT),
               font=F)

        elif kind == 'rect':
            fill = _resolve_color(el.get('fill', 'card'), st) if el.get('fill') else None
            line = _resolve_color(el.get('line', None), st) if el.get('line') else None
            shape = rc(slide, x, y, w, h, fill=fill, line=line,
                       rd=el.get('radius', False))
            if el.get('transparency'):
                _set_transparency(shape, el['transparency'])

        elif kind == 'oval':
            fill = _resolve_color(el.get('fill', 'accent'), st) if el.get('fill') else None
            shape = ov(slide, x, y, w, h, fill=fill,
                       transparency=el.get('transparency', 0))

        elif kind == 'image':
            img_path = el.get('path', '')
            if img_path and os.path.exists(img_path):
                slide.shapes.add_picture(img_path, I(x), I(y), I(w), I(h))

        elif kind == 'gradient':
            # Gradient rectangle: 2+ color stops, linear direction
            from pptx.oxml.ns import qn
            from lxml import etree
            s = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
            s.line.fill.background()
            gf = s.fill
            gf.gradient()
            colors = el.get('colors', ['navy', 'accent'])
            angle = el.get('angle', 0)
            gf.gradient_angle = angle
            stops = gf.gradient_stops
            # Set first two stops (always exist after gradient())
            stops[0].position = 0.0
            stops[0].color.rgb = _resolve_color(colors[0], st)
            stops[1].position = 1.0
            stops[1].color.rgb = _resolve_color(colors[1], st)
            # Add extra stops via XML for 3+ colors
            if len(colors) >= 3:
                gs_lst = s._element.find(qn('p:spPr')).find(qn('a:gradFill')).find(qn('a:gsLst'))
                for i, c in enumerate(colors[2:], 2):
                    pos = int((i / (len(colors) - 1)) * 100000)
                    gs = etree.SubElement(gs_lst, qn('a:gs'))
                    gs.set('pos', str(pos))
                    srgb = etree.SubElement(gs, qn('a:srgbClr'))
                    rgb = _resolve_color(c, st)
                    srgb.set('val', f'{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}')
                # Fix stop 1 position for even distribution
                stops[1].position = 1.0
            if el.get('transparency'):
                _set_transparency(s, el['transparency'])

        elif kind == 'line':
            color = _resolve_color(el.get('color', 'accent'), st)
            rc(slide, x, y, w, h, fill=color)


# ─────────────────────────────────────────────────────────────────────────────
# Routing tables
# ─────────────────────────────────────────────────────────────────────────────

CHROME_RENDERERS = {
    'cover':           render_cover,
    'agenda':          render_agenda,
    'section_divider': render_section_divider,
    'ending':          render_ending,
}

CONTENT_RENDERERS = {
    # Original 9
    'bullets_with_panel': render_bullets_with_panel,
    'two_column_check':   render_two_column_check,
    'cards_grid':         render_cards_grid,
    'criteria_rows':      render_criteria_rows,
    'scope_tiers':        render_scope_tiers,
    'two_panel':          render_two_panel,
    'two_column_steps':   render_two_column_steps,
    'scenario_cards':     render_scenario_cards,
    'survey':             render_survey,
    # New 5
    'stat_highlight':     render_stat_highlight,
    'timeline':           render_timeline,
    'quote_full':         render_quote_full,
    'center_focus':       render_center_focus,
    'comparison_table':   render_comparison_table,
    # New 2 + aliases
    'chart':              render_chart,
    'image_full':         render_image_full,
    'key_metrics':        render_stat_highlight,    # alias
    'quote_highlight':    render_quote_full,         # alias
    # Freeform
    'freeform':           render_freeform,
}


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description='EVYD PPT Generator — render content.json → PPTX')
    parser.add_argument('content', help='Path to content JSON file')
    parser.add_argument('--style', default=None,
                        help='Style preset name (default: from JSON meta.style, '
                             'then evyd_blue)')
    parser.add_argument('--output', '-o', default=None,
                        help='Output .pptx path (default: from JSON meta.output, '
                             'then output.pptx)')
    args = parser.parse_args()

    with open(args.content, encoding='utf-8') as f:
        content = json.load(f)
    meta = content.get('meta', {})

    style_name  = args.style  or meta.get('style',  'evyd_blue')
    output_path = args.output or meta.get('output', 'output.pptx')

    st = load_style(style_name)
    print(f'Style: {style_name}  →  {output_path}')

    prs = Presentation()
    prs.slide_width  = I(SW)
    prs.slide_height = I(SH)
    blank_layout = prs.slide_layouts[6]   # "Blank"

    slides  = content['slides']
    total   = len(slides)
    counter = 0

    for slide_data in slides:
        stype = slide_data.get('type', '')
        s = prs.slides.add_slide(blank_layout)

        if stype in CHROME_RENDERERS:
            CHROME_RENDERERS[stype](s, slide_data, st, counter, total)
        elif stype in CONTENT_RENDERERS:
            counter += 1
            _fill_bg(s, slide_data, st)
            CONTENT_RENDERERS[stype](s, slide_data, st, counter, total)
        else:
            print(f'  \u26a0  Unknown slide type "{stype}" \u2014 skipped')

    validate_and_fix(prs)
    prs.save(output_path)
    print(f'Slide count: {len(prs.slides)}')
    print(f'\u2713  Saved \u2192 {output_path}')


if __name__ == '__main__':
    main()
